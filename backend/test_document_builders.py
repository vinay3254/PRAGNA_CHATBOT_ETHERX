"""Test format builder functions produce valid, readable files."""
import os
import tempfile

from docx import Document as DocxReader
from openpyxl import load_workbook
from pptx import Presentation as PptxReader
from pypdf import PdfReader

from services.document_generator import _build_docx, _build_pdf, _build_pptx, _build_xlsx

STRUCTURE = {
    "title": "Test Document",
    "sections": [
        {"heading": "Overview", "bullets": ["Point one", "Point two"], "table": None},
        {"heading": "Numbers", "bullets": [], "table": [["Quarter", "Revenue"], ["Q1", "100"]]},
    ],
}

NO_TABLE_STRUCTURE = {
    "title": "Prose Only",
    "sections": [
        {"heading": "Intro", "bullets": ["First", "Second"], "table": None},
    ],
}


def test_docx():
    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "out.docx")
        _build_docx(STRUCTURE, path)
        assert os.path.getsize(path) > 0
        doc = DocxReader(path)
        text = "\n".join(p.text for p in doc.paragraphs)
        assert "Test Document" in text
        assert "Point one" in text
        assert len(doc.tables) == 1
        print("PASS: docx builder")


def test_pdf():
    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "out.pdf")
        _build_pdf(STRUCTURE, path)
        assert os.path.getsize(path) > 0
        reader = PdfReader(path)
        text = "".join(page.extract_text() or "" for page in reader.pages)
        assert "Test Document" in text
        assert "Point one" in text
        assert "Quarter" in text
        print("PASS: pdf builder")


def test_pptx():
    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "out.pptx")
        _build_pptx(STRUCTURE, path)
        assert os.path.getsize(path) > 0
        prs = PptxReader(path)
        assert len(prs.slides) == 3  # title slide + 2 section slides
        assert prs.slides[0].shapes.title.text == "Test Document"
        overview_body = "\n".join(p.text for p in prs.slides[1].placeholders[1].text_frame.paragraphs)
        assert "Point one" in overview_body
        numbers_body = "\n".join(p.text for p in prs.slides[2].placeholders[1].text_frame.paragraphs)
        assert "Quarter" in numbers_body and "Revenue" in numbers_body
        print("PASS: pptx builder")


def test_xlsx_with_table():
    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "out.xlsx")
        _build_xlsx(STRUCTURE, path)
        wb = load_workbook(path)
        assert "Numbers" in wb.sheetnames
        ws = wb["Numbers"]
        assert ws["A1"].value == "Quarter"
        print("PASS: xlsx builder with table")


def test_xlsx_fallback_no_table():
    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "out.xlsx")
        _build_xlsx(NO_TABLE_STRUCTURE, path)
        wb = load_workbook(path)
        assert "Summary" in wb.sheetnames
        ws = wb["Summary"]
        assert ws["A1"].value == "Section"
        assert ws["A2"].value == "Intro"
        print("PASS: xlsx builder fallback sheet")


if __name__ == "__main__":
    test_docx()
    test_pdf()
    test_pptx()
    test_xlsx_with_table()
    test_xlsx_fallback_no_table()
    print("All document builder tests passed.")
